#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Gera apresentação PowerPoint para clientes
Usando python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Definir cores
BLUE = RGBColor(37, 99, 235)
ORANGE = RGBColor(249, 115, 22)
GREEN = RGBColor(16, 185, 129)
PURPLE = RGBColor(139, 92, 246)
AMBER = RGBColor(245, 158, 11)
PINK = RGBColor(236, 72, 153)
CYAN = RGBColor(6, 182, 212)
INDIGO = RGBColor(99, 102, 241)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle, color):
    """Adiciona slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(72)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(40)
    subtitle_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, color):
    """Adiciona slide com conteúdo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Conteúdo
    y_position = 1.8
    for item in content_items:
        content_box = slide.shapes.add_textbox(Inches(1), Inches(y_position), Inches(8), Inches(0.8))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = item
        content_frame.paragraphs[0].font.size = Pt(28)
        content_frame.paragraphs[0].font.color.rgb = WHITE
        content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        y_position += 0.9
    
    return slide

# SLIDE 1: CAPA
print("Criando slide 1: Capa...")
add_title_slide(
    prs, 
    "⚽ Gestão de Clubes",
    "Sistema Completo para Gerenciamento de Times Amadores\nControle total: Jogadores • Jogos • Financeiro • Membros",
    BLUE
)

# SLIDE 2: O PROBLEMA
print("Criando slide 2: O Problema...")
add_content_slide(
    prs,
    "O Desafio",
    [
        "📋 Gerenciar mensalidades de forma desorganizada",
        "⚽ Controlar jogos e estatísticas manualmente",
        "💰 Dificuldade com controle financeiro",
        "👥 Sem comunicação centralizada com membros",
        "📊 Impossível gerar relatórios",
    ],
    ORANGE
)

# SLIDE 3: A SOLUÇÃO
print("Criando slide 3: A Solução...")
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = GREEN

# Título
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "A Solução"
title_frame.paragraphs[0].font.size = Pt(54)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = WHITE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Subtítulo
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.6))
sub_frame = sub_box.text_frame
sub_frame.text = "Uma plataforma web integrada que centraliza TUDO"
sub_frame.paragraphs[0].font.size = Pt(32)
sub_frame.paragraphs[0].font.color.rgb = WHITE
sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Features em 2 colunas
features = [
    ("👥 Gestão de Membros", "Cadastro, permissões e controle de acesso"),
    ("💳 Mensalidades", "Controle automatizado de pagamentos"),
    ("⚽ Jogos", "Agenda, resultados e estatísticas"),
    ("💰 Financeiro", "Entradas, saídas e relatórios"),
]

y = 2.3
for i, (title, desc) in enumerate(features):
    x = 0.8 if i % 2 == 0 else 5.5
    if i == 2:
        y = 4.5
    
    # Título feature
    feature_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4), Inches(0.5))
    feature_frame = feature_box.text_frame
    feature_frame.text = title
    feature_frame.paragraphs[0].font.size = Pt(22)
    feature_frame.paragraphs[0].font.bold = True
    feature_frame.paragraphs[0].font.color.rgb = WHITE
    
    # Descrição
    desc_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.5), Inches(4), Inches(0.7))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    desc_frame.text = desc
    desc_frame.paragraphs[0].font.size = Pt(18)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(220, 255, 220)

# SLIDE 4: FUNCIONALIDADES
print("Criando slide 4: Funcionalidades...")
add_content_slide(
    prs,
    "Funcionalidades Principais",
    [
        "🔐 Autenticação Segura - Login com email/senha",
        "📱 Responsivo 100% - Desktop, tablet e mobile",
        "⚡ Tempo Real - Dados sincronizam instantaneamente",
        "📊 Relatórios Automáticos - Gráficos e estatísticas",
    ],
    PURPLE
)

# SLIDE 5: CONTROLE DE ACESSO
print("Criando slide 5: Permissões...")
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = AMBER

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Controle de Permissões"
title_frame.paragraphs[0].font.size = Pt(54)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = WHITE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

roles = [
    "🔴 Admin Global - Acesso a tudo em todos clubes",
    "🟠 Administrador - Gestão completa do clube",
    "🟡 Diretor - Gerencia jogadores e jogos",
    "🟢 Jogador - Visualização apenas",
]

y = 2.0
for role in roles:
    role_box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(7), Inches(0.8))
    role_frame = role_box.text_frame
    role_frame.word_wrap = True
    role_frame.text = role
    role_frame.paragraphs[0].font.size = Pt(26)
    role_frame.paragraphs[0].font.color.rgb = WHITE
    role_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    y += 1.0

# SLIDE 6: DASHBOARD
print("Criando slide 6: Dashboard...")
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = PINK

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Dashboard Inteligente"
title_frame.paragraphs[0].font.size = Pt(54)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = WHITE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Simulação do dashboard
dashboard_items = [
    ("💵 Caixa Total", "R$ 2.500,00"),
    ("👥 Membros Ativos", "18"),
    ("💳 Mensalidades em Dia", "16"),
    ("⚽ Próximo Jogo", "Sábado às 15h"),
]

y = 2.0
for label, value in dashboard_items:
    item_box = slide.shapes.add_textbox(Inches(2), Inches(y), Inches(6), Inches(0.9))
    item_frame = item_box.text_frame
    item_frame.word_wrap = True
    
    p = item_frame.paragraphs[0]
    p.text = f"{label}: {value}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    y += 1.0

# SLIDE 7: BENEFÍCIOS
print("Criando slide 7: Benefícios...")
add_content_slide(
    prs,
    "Por Que Contratar?",
    [
        "💰 Reduz custos com administração manual",
        "⏱️ Economiza tempo em processos repetitivos",
        "📈 Aumenta transparência com membros",
        "🔒 Protege dados com segurança em nuvem",
        "📊 Facilita decisões com relatórios automáticos",
        "📱 Acessível anywhere - qualquer dispositivo",
    ],
    CYAN
)

# SLIDE 8: TECNOLOGIA
print("Criando slide 8: Tecnologia...")
add_content_slide(
    prs,
    "Arquitetura Técnica",
    [
        "☁️ Infraestrutura em Nuvem - Firebase 99.95% uptime",
        "⚡ Escalável - Suporta 1 a 10.000+ usuários",
        "🔐 Ultra Seguro - Criptografia e backups automáticos",
        "🌍 Suporte 24/7 - Infraestrutura gerenciada",
    ],
    INDIGO
)

# SLIDE FINAL: CHAMADA
print("Criando slide 9: Chamada...")
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BLUE

# Título
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_frame.text = "Pronto para Começar?"
title_frame.paragraphs[0].font.size = Pt(64)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = WHITE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Subtítulo
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Transforme a gestão do seu clube hoje!"
subtitle_frame.paragraphs[0].font.size = Pt(40)
subtitle_frame.paragraphs[0].font.color.rgb = WHITE
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Contato
contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
contact_frame = contact_box.text_frame
contact_frame.word_wrap = True

p1 = contact_frame.paragraphs[0]
p1.text = "📧 william.santos@company.com"
p1.font.size = Pt(32)
p1.font.color.rgb = WHITE
p1.alignment = PP_ALIGN.CENTER

p2 = contact_frame.add_paragraph()
p2.text = "📱 (11) 99999-9999"
p2.font.size = Pt(32)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER

# Salvar
output_file = "apresentacao-clientes.pptx"
prs.save(output_file)
print(f"✅ Apresentação criada com sucesso: {output_file}")
